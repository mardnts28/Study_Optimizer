import re
import io
import PyPDF2
import docx
from pptx import Presentation
from datetime import date, timedelta
from django.db.models import Sum, Count
from django.contrib.postgres.search import SearchVector, SearchQuery, SearchRank
from .models import Task, SharedMaterial, Comment, SummarizedDocument, ScheduleItem

def extract_text_from_file(uploaded_file):
    """
    Extracts text from various file formats (PDF, DOCX, PPTX, TXT).
    """
    file_name = uploaded_file.name
    content = ""
    
    if file_name.lower().endswith('.pdf'):
        reader = PyPDF2.PdfReader(uploaded_file)
        for page in reader.pages[:20]: 
            text = page.extract_text()
            if text: content += text + " "
    
    elif file_name.lower().endswith('.docx'):
        doc = docx.Document(uploaded_file)
        content = " ".join([p.text for p in doc.paragraphs[:100]])
        
    elif file_name.lower().endswith('.pptx'):
        prs = Presentation(uploaded_file)
        for slide in prs.slides[:30]:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    content = content + str(shape.text) + " "
                    
    elif file_name.lower().endswith('.txt'):
        content = uploaded_file.read().decode('utf-8', errors='ignore')
        
    return content.strip()

def highlight_keywords(text, important_words):
    """
    Highlights important words in text using markdown bold.
    """
    for word in important_words:
        if len(word) > 3:
            pattern = rf'\b({re.escape(word)})\b'
            text = re.sub(pattern, r'**\1**', text, flags=re.IGNORECASE)
    return text

def generate_document_summary(content, file_name):
    """
    Generates a polished, structured summary with a 'Name - Description' format.
    No bolding used, as per user request.
    """
    content_clean = re.sub(r'\s+', ' ', content).strip()
    raw_sents = re.split(r'(?<=[.!?])\s+', content_clean)
    sentences = [s.strip() for s in raw_sents if len(s.strip()) > 20]
    
    words = re.findall(r'\w+', content_clean.lower())
    stopwords = set(['the', 'and', 'is', 'in', 'to', 'of', 'a', 'that', 'it', 'for', 'on', 'with', 'as', 'this', 'was', 'at', 'by', 'an', 'be', 'are', 'which', 'from', 'or', 'their', 'we', 'your', 'has', 'have', 'were', 'not', 'can', 'will', 'but', 'all', 'they', 'he', 'she', 'his', 'her', 'who', 'about', 'some', 'more', 'so', 'one', 'out', 'up', 'down', 'into', 'over', 'after', 'before', 'then', 'once', 'just', 'only', 'than', 'them', 'if', 'there', 'when', 'any', 'each', 'other', 'been', 'would', 'could', 'should'])
    
    freq_dict = {}
    for w in words:
        if w not in stopwords and len(w) > 4:
            freq_dict[w] = freq_dict.get(w, 0) + 1
    
    score_map = {}
    for i, sent in enumerate(sentences):
        s_words = re.findall(r'\w+', sent.lower())
        val = sum(freq_dict.get(sw, 0) for sw in s_words)
        if len(s_words) > 8:
            score_map[i] = val / (len(s_words) ** 0.5)

    # Scale summary length with document size: 10 sentences min, up to 25 for long docs
    summary_size = min(max(10, len(sentences) // 20), 25)
    top_indices = sorted(list(score_map.keys()), key=lambda x: score_map[x], reverse=True)[:summary_size]
    top_indices.sort()
    relevant_sentences = [sentences[idx] for idx in top_indices]

    def format_line(s, prefix):
        words = s.split()
        if len(words) > 5:
            # Topic - Description format
            topic = " ".join(words[:2]).strip(',.:; ').upper()
            description = " ".join(words[2:]).strip()
            return f"{prefix} {topic}: {description}"
        return f"{prefix} {s}"

    header = f"DOCUMENT SUMMARY: {file_name.upper()}"
    highlights = [format_line(s, "✅") for s in relevant_sentences[:3]]
    overview = " ".join(relevant_sentences[:2]) if relevant_sentences else "Minimal content found."
    details = [format_line(s, "•") for s in relevant_sentences[1:min(len(relevant_sentences), 15)]]
    
    summary_parts = [
        header,
        "",
        "KEY HIGHLIGHTS",
        "\n".join(highlights),
        "",
        "OVERVIEW",
        overview,
        "",
        "CORE CONCEPTS",
        "\n".join(details),
        "",
        "TAKEAWAY",
        "This resource provides a clear technical foundation. Reviewing the points above will ensure a strong grasp of the primary subject matter."
    ]
    
    final_output = "\n".join(summary_parts)
    return final_output, f"Analysis of {file_name}"

def generate_batch_synthesis(doc_ids, user):
    """
    Synthesizes multiple documents into a single thematic summary.
    """
    docs = SummarizedDocument.objects.filter(id__in=doc_ids, user=user)
    if not docs.exists():
        return None
        
    all_text_lower = " ".join([str(d.summary_text) for d in docs]).lower()
    
    doc_categories = []
    if any(k in all_text_lower for k in ['vpn', 'network', 'protocol']): doc_categories.append('Network Security')
    if any(k in all_text_lower for k in ['iam', 'policy', 'access', 'auth']): doc_categories.append('Identity Management')
    if any(k in all_text_lower for k in ['cloud', 'aws', 'azure', 'serverless']): doc_categories.append('Cloud Infrastructure')
    if any(k in all_text_lower for k in ['setup', 'config', 'deploy']): doc_categories.append('System Configuration')
    
    unique_categories = list(set(doc_categories))
    category_context = " & ".join(unique_categories) if unique_categories else "Technical Systems"

    batch_output = f"📊 BATCH SYNTHESIS: {category_context}\n\n"
    
    if len(unique_categories) >= 2:
        batch_output += f"🔄 Layered Insights:\n"
        batch_output += f"These documents define a cohesive {category_context} strategy. For example, {unique_categories[0]} provides the perimeter security necessary for {unique_categories[1]} to function effectively. This horizontal integration demonstrates that technical success is dependent on cross-layered configuration alignment.\n\n"
    else:
        batch_output += f"🔄 Strategic Focus:\n"
        batch_output += f"The combined resources provide a comprehensive framework for {category_context}. They establish detailed standards for deployment, auditing, and optimization, ensuring that every configuration step is backed by technical best practices.\n\n"
    
    final_sponsor = f"🎯 Final Sponsor Takeaway: This curated set establishes the exact critical competencies required for modern professional excellence. Mastering these overlapping themes bridges the gap between individual configurations and strategic infrastructure management."
    batch_output += final_sponsor
    
    return batch_output

def calculate_user_metrics(user):
    """
    Calculates all student metrics: streak, study hours, tasks, levels, and progress.
    """
    user_tasks = Task.objects.filter(user=user)
    total_tasks = user_tasks.count()
    completed_tasks = user_tasks.filter(completed=True).count()
    summaries = SummarizedDocument.objects.filter(user=user)
    
    # Calculate Completion Rate
    completion_rate = round((completed_tasks / total_tasks * 100), 1) if total_tasks > 0 else 0
    
    # Calculate Streak
    activity_dates = set()
    for t in user_tasks.filter(completed=True):
        activity_dates.add(t.created_at.date())
    for s in summaries:
        activity_dates.add(s.created_at.date())
    
    sorted_dates = sorted(list(activity_dates), reverse=True)
    streak = 0
    today = date.today()
    if sorted_dates:
        if sorted_dates[0] == today or sorted_dates[0] == today - timedelta(days=1):
            current_date = sorted_dates[0]
            streak = 1
            for i in range(1, len(sorted_dates)):
                if sorted_dates[i] == current_date - timedelta(days=1):
                    streak += 1
                    current_date = sorted_dates[i]
                else: break
    
    # Calculate Study Hours (2h/task, 1h/summary/share)
    total_shares = SharedMaterial.objects.filter(author=user).count()
    study_hours = (completed_tasks * 2) + ((summaries.count() + total_shares) * 1)
    
    # Calculate Level (1 level every 5 completed tasks)
    user_level = (completed_tasks // 5) + 1
    next_level_progress = int(((completed_tasks % 5) / 5) * 100)
    
    # Subject Distribution
    subjects = {}
    for t in user_tasks:
        sub = t.subject if t.subject else 'General'
        subjects[sub] = subjects.get(sub, 0) + 1
    for s in summaries:
        sub = s.subject if s.subject else 'General'
        subjects[sub] = subjects.get(sub, 0) + 1
    
    subject_labels = list(subjects.keys())
    subject_data = list(subjects.values())

    # Weekly hours (last 4 weeks)
    weekly_trend = []
    for i in range(3, -1, -1):
        start_date = today - timedelta(days=(i+1)*7)
        end_date = today - timedelta(days=i*7)
        tasks_in_week = user_tasks.filter(completed=True, created_at__range=(start_date, end_date)).count()
        sums_in_week = summaries.filter(created_at__range=(start_date, end_date)).count()
        weekly_trend.append((tasks_in_week * 2) + sums_in_week)

    return {
        'total_tasks': total_tasks,
        'completed_count': completed_tasks,
        'completion_rate': completion_rate,
        'streak': streak,
        'study_hours': study_hours,
        'summaries_count': summaries.count(),
        'user_level': user_level,
        'next_level_progress': next_level_progress,
        'docs_count': SharedMaterial.objects.filter(author=user).count() + summaries.count(),
        'subject_labels': subject_labels,
        'subject_data': subject_data,
        'weekly_hours_trend': weekly_trend,
    }

def search_summarized_documents(user, query):
    """
    Performs full-text search on summarized documents.
    """
    if not query:
        return []
    
    vector = SearchVector('file_name', weight='A') + SearchVector('summary_text', weight='B')
    search_query = SearchQuery(query)
    
    return SummarizedDocument.objects.filter(user=user).annotate(
        rank=SearchRank(vector, search_query)
    ).filter(rank__gte=0.1).order_by('-rank')
