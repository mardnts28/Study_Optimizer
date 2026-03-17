from django.shortcuts import render, redirect, get_object_or_404
from django.contrib.auth.models import User
from django.contrib import messages
from django.contrib.auth import authenticate, login, logout
from django.contrib.auth.decorators import login_required
from django.views.decorators.csrf import csrf_protect
from django.views.decorators.http import require_POST
from django.http import JsonResponse
from django.db.models import Sum, Count
import json
from .models import Task, SharedMaterial, Comment, SummarizedDocument, ScheduleItem
from datetime import date, timedelta

@csrf_protect
def index(request):
    return render(request, "main/index.html")

@csrf_protect
def register(request):
    if request.user.is_authenticated:
        return redirect("dashboard")
    if request.method == "POST":
        username = request.POST.get('username', '').strip()
        email = request.POST.get('email', '').strip()
        password = request.POST.get('password', '')
        password2 = request.POST.get('password2', '')

        if password != password2:
            messages.error(request, "Passwords do not match.")
            return redirect('register')

        if User.objects.filter(username=username).exists():
            messages.error(request, "Username already taken.")
            return redirect('register')

        if User.objects.filter(email=email).exists():
            messages.error(request, "Email is already registered.")
            return redirect('register')

        User.objects.create_user(username=username, email=email, password=password)
        messages.success(request, "Account created successfully. Please log in.")
        return redirect('login')

    return render(request, 'main/register.html')

def login_view(request):
    if request.user.is_authenticated:
        return redirect("dashboard")
    if request.method == "POST":
        email = request.POST.get("email")
        password = request.POST.get("password")

        try:
            user_obj = User.objects.get(email=email)
            user = authenticate(request, username=user_obj.username, password=password)
        except User.DoesNotExist:
            user = None

        if user is not None:
            login(request, user)
            return redirect("dashboard")  
        else:
            messages.error(request, "Invalid email or password.")

    return render(request, "main/login.html")

def logout_view(request):
    logout(request)
    return redirect("login")

@login_required
@csrf_protect
def collaborate(request):
    materials = SharedMaterial.objects.all().order_by('-created_at')
    materials_list = []
    
    total_community_likes = 0
    for m in materials:
        m_likes = m.likes.count()
        total_community_likes += m_likes
        materials_list.append({
            'id': m.id,
            'title': m.title,
            'author': m.author.username,
            'authorInitials': m.author.username[:2].upper(),
            'authorColor': '#8C1007' if m.author == request.user else '#4B5563',
            'period': m.period,
            'subject': m.subject,
            'preview': m.content,
            'likes': m_likes,
            'views': m.views,
            'comments': m.comments.count(),
            'timeAgo': 'Just now',
            'emoji': m.emoji,
            'liked': m.likes.filter(id=request.user.id).exists(),
            'tags': [m.subject]
        })
    
    active_students_count = User.objects.count()
    
    return render(request, "main/collaborate.html", {
        'materials_json': json.dumps(materials_list),
        'active_students': active_students_count,
        'total_community_likes': total_community_likes
    })

@login_required
@require_POST
def share_material(request):
    try:
        data = json.loads(request.body)
        material = SharedMaterial.objects.create(
            author=request.user,
            title=data.get('title'),
            subject=data.get('subject'),
            period=data.get('period'),
            content=data.get('preview'),
            emoji='📄'
        )
        return JsonResponse({
            'status': 'success',
            'material': {
                'id': material.id,
                'title': material.title,
                'author': material.author.username,
                'authorInitials': material.author.username[:2].upper(),
                'authorColor': '#8C1007',
                'period': material.period,
                'subject': material.subject,
                'preview': material.content,
                'likes': 0,
                'views': 0,
                'comments': 0,
                'timeAgo': 'Just now',
                'emoji': material.emoji,
                'liked': False,
                'tags': [material.subject]
            }
        })
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

@login_required
@require_POST
def toggle_like_material(request, material_id):
    material = get_object_or_404(SharedMaterial, id=material_id)
    if material.likes.filter(id=request.user.id).exists():
        material.likes.remove(request.user)
        liked = False
    else:
        material.likes.add(request.user)
        liked = True
    return JsonResponse({'status': 'success', 'liked': liked, 'likes_count': material.likes.count()})

@login_required
def get_material_comments(request, material_id):
    material = get_object_or_404(SharedMaterial, id=material_id)
    comments = material.comments.all().order_by('-created_at')
    comments_list = [{
        'id': c.id,
        'author': c.author.username,
        'authorInitials': c.author.username[:2].upper(),
        'authorColor': '#8C1007' if c.author == request.user else '#4B5563',
        'text': c.text,
        'timeAgo': 'Just now'
    } for c in comments]
    return JsonResponse({'status': 'success', 'comments': comments_list})

@login_required
@require_POST
def add_comment(request, material_id):
    material = get_object_or_404(SharedMaterial, id=material_id)
    data = json.loads(request.body)
    comment = Comment.objects.create(
        material=material,
        author=request.user,
        text=data.get('text')
    )
    return JsonResponse({
        'status': 'success',
        'comment': {
            'id': comment.id,
            'author': comment.author.username,
            'authorInitials': comment.author.username[:2].upper(),
            'authorColor': '#8C1007',
            'text': comment.text,
            'timeAgo': 'Just now'
        }
    })

@login_required
@csrf_protect
def dashboard(request):
    user_tasks = Task.objects.filter(user=request.user)
    active_tasks = user_tasks.filter(completed=False).order_by('due_date')
    
    total_tasks = user_tasks.count()
    completed_tasks = user_tasks.filter(completed=True).count()
    completion_rate = round((completed_tasks / total_tasks * 100), 1) if total_tasks > 0 else 0
    
    upcoming_tasks_list = []
    for t in active_tasks[:4]:
        days_left = (t.due_date - date.today()).days
        upcoming_tasks_list.append({
            'title': t.title,
            'date': t.due_date.strftime('%b %d'),
            'priority': t.priority,
            'period': 'General',
            'daysLeft': days_left if days_left >= 0 else 0
        })

    # General task progress instead of period breakdown
    total_active = active_tasks.count()
    completed_today = user_tasks.filter(completed=True, created_at__date=date.today()).count()

    # Weekly Study Hours (Mon-Sun)
    today = date.today()
    start_of_week = today - timedelta(days=today.weekday())
    daily_hours = []
    for i in range(7):
        day = start_of_week + timedelta(days=i)
        tasks_on_day = user_tasks.filter(completed=True, created_at__date=day).count()
        summaries_on_day = SummarizedDocument.objects.filter(user=request.user, created_at__date=day).count()
        daily_hours.append((tasks_on_day * 2) + (summaries_on_day * 1))

    context = {
        'total_tasks': total_tasks,
        'completed_count': completed_tasks,
        'completion_rate': completion_rate,
        'upcoming_tasks_json': json.dumps(upcoming_tasks_list),
        'docs_count': SharedMaterial.objects.filter(author=request.user).count() + SummarizedDocument.objects.filter(user=request.user).count(),
        'study_hours': (completed_tasks * 2) + (SummarizedDocument.objects.filter(user=request.user).count() * 1),
        'weekly_hours_list': json.dumps(daily_hours),
        'schedule_items_json': json.dumps([{
            'id': item.id,
            'day': item.day,
            'time': item.time,
            'activity': item.activity,
            'color': item.color
        } for item in ScheduleItem.objects.filter(user=request.user)]),
    }
    return render(request, "main/dashboard.html", context)

@login_required
@csrf_protect
def progress(request):
    user_tasks = Task.objects.filter(user=request.user)
    total_tasks = user_tasks.count()
    completed_tasks = user_tasks.filter(completed=True).count()
    summaries = SummarizedDocument.objects.filter(user=request.user)
    
    # Calculate Streak
    activity_dates = set()
    # Note: Using created_at as completion date proxy if no explicit completion timestamp exists
    for t in user_tasks.filter(completed=True):
        activity_dates.add(t.created_at.date())
    for s in summaries:
        activity_dates.add(s.created_at.date())
    
    sorted_dates = sorted(list(activity_dates), reverse=True)
    streak = 0
    today = date.today()
    
    if sorted_dates:
        # Check if the streak is still alive (active today or yesterday)
        if sorted_dates[0] == today or sorted_dates[0] == today - timedelta(days=1):
            current_date = sorted_dates[0]
            streak = 1
            for i in range(1, len(sorted_dates)):
                if sorted_dates[i] == current_date - timedelta(days=1):
                    streak += 1
                    current_date = sorted_dates[i]
                else:
                    break
    
    # Estimate Study Hours: 2h per completed task, 1h per summary
    study_hours_total = (completed_tasks * 2) + (summaries.count() * 1)
    
    # Subject distribution
    subjects = {}
    for t in user_tasks:
        sub = t.subject if t.subject else 'General'
        subjects[sub] = subjects.get(sub, 0) + 1
    for s in summaries:
        sub = s.subject if s.subject else 'General'
        subjects[sub] = subjects.get(sub, 0) + 1
    
    subject_labels = list(subjects.keys())
    subject_data = list(subjects.values())

    # Weekly hours (last 4 weeks as requested by chart labels 'Week 1-4' originally, 
    # but let's do last 4 segments of activity)
    # Actually the chart has labels 'Week 1', 'Week 2', 'Week 3', 'Week 4'
    weekly_trend = []
    for i in range(3, -1, -1):
        start_date = today - timedelta(days=(i+1)*7)
        end_date = today - timedelta(days=i*7)
        tasks_in_week = user_tasks.filter(completed=True, created_at__range=(start_date, end_date)).count()
        sums_in_week = summaries.filter(created_at__range=(start_date, end_date)).count()
        weekly_trend.append((tasks_in_week * 2) + sums_in_week)

    period_stats = [
        {'period': 'General Progress', 'completed': completed_tasks, 'total': total_tasks}
    ]
    
    context = {
        'total_tasks': total_tasks,
        'completed_count': completed_tasks,
        'completion_rate': round((completed_tasks / total_tasks * 100), 1) if total_tasks > 0 else 0,
        'streak': streak,
        'study_hours': study_hours_total,
        'summaries_count': summaries.count(),
        'period_stats_json': json.dumps(period_stats),
        'subject_labels_json': json.dumps(subject_labels),
        'subject_data_json': json.dumps(subject_data),
        'weekly_hours_json': json.dumps(weekly_trend),
    }
    return render(request, "main/progress.html", context)

@login_required
@csrf_protect
def upload(request):
    recent_summaries = SummarizedDocument.objects.filter(user=request.user).order_by('-created_at')[:10]
    summaries_list = []
    for s in recent_summaries:
        summaries_list.append({
            'id': s.id,
            'title': s.file_name,
            'period': 'General',
            'date': s.created_at.strftime('%b %d'),
            'emoji': s.emoji,
            'summary': s.summary_text
        })
    return render(request, "main/upload.html", {
        'recent_summaries_json': json.dumps(summaries_list)
    })

import PyPDF2
import io

@login_required
@require_POST
def summarize_doc(request):
    try:
        if 'file' not in request.FILES:
            return JsonResponse({'status': 'error', 'message': 'No file uploaded'}, status=400)
        
        uploaded_file = request.FILES['file']
        period = 'General'
        file_name = uploaded_file.name
        content = ""
        
        # Text Extraction per File Type
        if file_name.lower().endswith('.pdf'):
            import PyPDF2
            reader = PyPDF2.PdfReader(uploaded_file)
            for page in reader.pages[:15]: # Increased page limit
                text = page.extract_text()
                if text: content += text + " "
        
        elif file_name.lower().endswith('.docx'):
            import docx
            doc = docx.Document(uploaded_file)
            content = " ".join([p.text for p in doc.paragraphs[:50]]) # Limit paragraphs
            
        elif file_name.lower().endswith('.pptx'):
            from pptx import Presentation
            prs = Presentation(uploaded_file)
            for slide in prs.slides[:20]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        content = content + str(shape.text) + " "
                        
        elif file_name.lower().endswith('.txt'):
            content = uploaded_file.read().decode('utf-8')
        else:
            return JsonResponse({'status': 'error', 'message': 'Unsupported file type. Please use PDF, DOCX, PPTX, or TXT.'}, status=400)

        if not content.strip():
            return JsonResponse({'status': 'error', 'message': 'Could not extract text from document'}, status=400)

        # ---------------------------------------------------------
        # AI Summarization Logic (Refined)
        # ---------------------------------------------------------
        import re
        from collections import Counter

        # Cleanup and tokenize
        content_clean = re.sub(r'\s+', ' ', content).strip()
        raw_sents = re.split(r'(?<=[.!?])\s+', content_clean)
        sentences = []
        for s in raw_sents:
            s_clean = s.strip()
            if len(s_clean) > 10:
                sentences.append(s_clean)
        
        # Frequency analysis (Standard dict for linter stability)
        words = re.findall(r'\w+', content_clean.lower())
        stopwords = set(['the', 'and', 'is', 'in', 'to', 'of', 'a', 'that', 'it', 'for', 'on', 'with', 'as', 'this', 'was', 'at', 'by', 'an', 'be', 'are', 'which', 'from', 'or', 'their', 'we', 'your', 'has', 'have', 'were', 'not', 'can', 'will', 'but', 'all', 'they', 'he', 'she', 'his', 'her', 'who', 'about', 'some', 'more', 'so', 'one', 'out', 'up', 'down', 'into', 'over', 'after', 'before', 'then', 'once', 'just', 'only', 'than', 'them', 'if', 'there', 'when', 'any', 'each', 'other', 'been', 'would', 'could', 'should'])
        
        freq_dict = {}
        for w in words:
            if w not in stopwords and len(w) > 3:
                freq_dict[w] = freq_dict.get(w, 0) + 1
        
        # Technical-Aware Scoring
        score_map = {}
        tech_keys = ['config', 'setup', 'policy', 'cli', 'network', 'vpn', 'iam', 's3', 'aws', 'router', 'switch', 'encryption', 'secure']
        for i in range(len(sentences)):
            sent = sentences[i]
            s_lower = sent.lower()
            s_words = re.findall(r'\w+', s_lower)
            
            val = 0.0
            for sw in s_words:
                val += float(freq_dict.get(sw, 0))
            
            t_weight = 1.0
            if any(tk in s_lower for tk in tech_keys): t_weight *= 1.5
            if re.search(r'\d+\.\d+\.\d+\.\d+', sent): t_weight *= 2.0
            if re.search(r'[A-Za-z0-9_-]+/[A-Za-z0-9_-]+', sent): t_weight *= 1.2
            
            total_val = val * t_weight
            if len(s_words) > 5:
                score_map[i] = total_val / (float(len(s_words)) ** 0.5)        # Extraction logic
        map_keys = list(score_map.keys())
        top_indices = sorted(map_keys, key=lambda x: score_map[x], reverse=True)
        if len(top_indices) > 10:
            top_indices = top_indices[:10]
        top_indices.sort()
        
        relevant_sentences = []
        for idx in top_indices:
            relevant_sentences.append(sentences[idx])

        # Technical Synthesis Workflow 
        sorted_freq = sorted(freq_dict.items(), key=lambda x: x[1], reverse=True)
        top_context_words = []
        for wf in sorted_freq:
            if len(wf[0]) > 4:
                top_context_words.append(str(wf[0]).capitalize())
            if len(top_context_words) >= 2: break
        
        tech_context = " & ".join(top_context_words) if top_context_words else "Infrastructure"
        title_line = f"📘 {tech_context} – Technical Infrastructure Overview"

        # Executive Highlights
        highlights = []
        icons = ["✅", "⚠️", "🚧", "💡"]
        action_keywords = ['configure', 'assigned', 'verified', 'created', 'enabled', 'tested', 'implemented', 'policy']
        
        high_pool = []
        for s in relevant_sentences:
            if any(ak in s.lower() for ak in action_keywords):
                high_pool.append(s)
        for s in relevant_sentences:
            if s not in high_pool:
                high_pool.append(s)
        
        num_highlights = min(4, len(high_pool))
        for i in range(num_highlights):
            highlights.append(f"{icons[i]} {high_pool[i]}")

        # Breakdown Sections
        if len(relevant_sentences) > 0:
            overview_text = relevant_sentences[0]
        else:
            overview_text = "Technical configuration analysis."
            
        tech_facts = []
        for s in relevant_sentences:
            if re.search(r'\d', s):
                tech_facts.append(s)
        
        # Key Details 
        detail_pool = []
        if len(tech_facts) > 0:
            detail_pool = tech_facts
        elif len(relevant_sentences) > 1:
            for i in range(1, min(6, len(relevant_sentences))):
                detail_pool.append(relevant_sentences[i])
        
        detail_lines = []
        for ds in detail_pool:
            detail_lines.append(f"• {ds}")
        
        key_details_text = "\n".join(detail_lines) if detail_lines else "• Verified system integrity parameters."
        
        # Implications based on content flags
        c_low = content_clean.lower()
        if "vpn" in c_low or "ips" in c_low:
            next_step = "Practice tunnel verification and encryption protocol analysis to ensure network-level isolation."
        elif "iam" in c_low or "policy" in c_low:
            next_step = "Review least-privilege principles and audit user permission boundaries to prevent privilege escalation."
        else:
            next_step = f"Deepen understanding of {tech_context} configurations to optimize system security and performance."

        takeaway_line = f"🎯 Takeaway: Mastering these {tech_context} configurations builds a professional foundation for your success."

        # Final Formatting (Sponsor-Ready Synthesis)
        final_summary = f"{title_line}\n\n"
        final_summary += "🔑 Executive Highlights\n" + "\n".join(highlights) + "\n\n"
        final_summary += "📂 Three-Part Breakdown\n"
        final_summary += f"Overview: {overview_text}\n\n"
        final_summary += f"Key Details:\n{key_details_text}\n\n"
        final_summary += f"Implications / Next Steps: {next_step}\n\n"
        final_summary += takeaway_line

        # Save to DB
        doc = SummarizedDocument.objects.create(
            user=request.user,
            file_name=file_name,
            period=period,
            summary_text=final_summary,
            emoji='📄'
        )

        return JsonResponse({
            'status': 'success',
            'summary': final_summary,
            'title': title_line,
            'doc_id': doc.id
        })
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=500)

@login_required
@require_POST
def summarize_batch(request):
    try:
        data = json.loads(request.body)
        doc_ids = data.get('doc_ids', [])
        period = 'General'
        
        if not doc_ids:
            return JsonResponse({'status': 'error', 'message': 'No documents provided for batch processing'}, status=400)
            
        docs = SummarizedDocument.objects.filter(id__in=doc_ids, user=request.user)
        if not docs.exists():
            return JsonResponse({'status': 'error', 'message': 'Documents not found'}, status=404)
            
        # Analysis for patterns & Layered Insights
        all_text_lower = " ".join([str(d.summary_text) for d in docs]).lower()
        
        doc_categories = []
        if 'vpn' in all_text_lower or 'network' in all_text_lower: doc_categories.append('Network Security')
        if 'iam' in all_text_lower or 'policy' in all_text_lower: doc_categories.append('Identity Management')
        if 'cloud' in all_text_lower or 'aws' in all_text_lower: doc_categories.append('Cloud Infrastructure')
        if 'setup' in all_text_lower or 'config' in all_text_lower: doc_categories.append('System Configuration')
        
        unique_categories = list(set(doc_categories))
        category_context = " & ".join(unique_categories) if unique_categories else "Technical Systems"

        # Construct Combined Insights (Synthesized Pattern Recognition)
        batch_output = f"📊 BATCH SYNTHESIS: {category_context}\n\n"
        
        if len(unique_categories) >= 2:
            batch_output += f"🔄 Layered Insights:\n"
            batch_output += f"These documents define a cohesive {category_context} strategy. For example, {unique_categories[0]} provides the perimeter security necessary for {unique_categories[1]} to function effectively. This vertical integration demonstrates that technical success is dependent on cross-layered configuration alignment.\n\n"
        else:
            batch_output += f"🔄 Strategic Focus:\n"
            batch_output += f"The combined resources provide a comprehensive framework for {category_context}. They establish detailed standards for deployment, auditing, and optimization, ensuring that every configuration step is backed by technical best practices.\n\n"
        
        final_sponsor = f"🎯 Final Sponsor Takeaway: This curated set establishes the exact critical competencies required for modern professional excellence. Mastering these overlapping themes bridges the gap between individual configurations and strategic infrastructure management."

        batch_output += final_sponsor
        
        return JsonResponse({
            'status': 'success',
            'combined_summary': batch_output
        })
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=500)

@login_required
@csrf_protect
def tasks(request):
    user_tasks = Task.objects.filter(user=request.user).order_by('completed', 'due_date')
    tasks_json = []
    for t in user_tasks:
        tasks_json.append({
            'id': t.id,
            'title': t.title,
            'period': t.period,
            'priority': t.priority,
            'dueDate': t.due_date.strftime('%Y-%m-%d'),
            'completed': t.completed
        })
    
    return render(request, "main/tasks.html", {'tasks_data': json.dumps(tasks_json)})

@login_required
@require_POST
def add_task(request):
    try:
        data = json.loads(request.body)
        task = Task.objects.create(
            user=request.user,
            title=data.get('title'),
            period=data.get('period'),
            priority=data.get('priority'),
            due_date=data.get('dueDate')
        )
        return JsonResponse({
            'status': 'success',
            'task': {
                'id': task.id,
                'title': task.title,
                'period': task.period,
                'priority': task.priority,
                'dueDate': task.due_date.strftime('%Y-%m-%d'),
                'completed': task.completed
            }
        })
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

@login_required
@require_POST
def edit_task(request, task_id):
    try:
        task = get_object_or_404(Task, id=task_id, user=request.user)
        data = json.loads(request.body)
        task.title = data.get('title', task.title)
        task.period = 'General'
        task.priority = data.get('priority', task.priority)
        task.due_date = data.get('dueDate', task.due_date)
        task.save()
        return JsonResponse({
            'status': 'success',
            'task': {
                'id': task.id,
                'title': task.title,
                'period': task.period,
                'priority': task.priority,
                'dueDate': task.due_date.strftime('%Y-%m-%d'),
                'completed': task.completed
            }
        })
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

@login_required
@require_POST
def toggle_task(request, task_id):
    task = get_object_or_404(Task, id=task_id, user=request.user)
    task.completed = not task.completed
    task.save()
    return JsonResponse({'status': 'success', 'completed': task.completed})

@login_required
@require_POST
def delete_task(request, task_id):
    task = get_object_or_404(Task, id=task_id, user=request.user)
    task.delete()
    return JsonResponse({'status': 'success'})

@login_required
@csrf_protect
@login_required
@csrf_protect
def profile(request):
    user_tasks = Task.objects.filter(user=request.user)
    total_tasks = user_tasks.count()
    completed_tasks = user_tasks.filter(completed=True).count()
    completion_rate = round((completed_tasks / total_tasks * 100), 1) if total_tasks > 0 else 0
    summaries_count = SummarizedDocument.objects.filter(user=request.user).count()
    study_hours = (completed_tasks * 2) + (summaries_count * 1)
    
    # Calculate Level: 1 level every 5 completed tasks
    user_level = (completed_tasks // 5) + 1
    next_level_progress = int(((completed_tasks % 5) / 5) * 100)

    context = {
        'user_level': user_level,
        'next_level_progress': next_level_progress,
        'docs_count': summaries_count,
        'completed_count': completed_tasks,
        'total_tasks': total_tasks,
        'completion_rate': completion_rate,
        'study_hours': study_hours,
    }
    return render(request, "main/profile.html", context)

@login_required
@require_POST
def add_schedule_item(request):
    try:
        data = json.loads(request.body)
        item = ScheduleItem.objects.create(
            user=request.user,
            day=data.get('day'),
            time=data.get('time'),
            activity=data.get('activity'),
            color=data.get('color', 'blue')
        )
        return JsonResponse({
            'status': 'success',
            'item': {
                'id': item.id,
                'day': item.day,
                'time': item.time,
                'activity': item.activity,
                'color': item.color
            }
        })
    except Exception as e:
        return JsonResponse({'status': 'error', 'message': str(e)}, status=400)

@login_required
@require_POST
def delete_schedule_item(request, item_id):
    item = get_object_or_404(ScheduleItem, id=item_id, user=request.user)
    item.delete()
    return JsonResponse({'status': 'success'})
